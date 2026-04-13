"""
pipeline/summarizer.py — Text extraction and summarization pipeline.

Three-level rolling summarization:
  Level 1 — Document-group summary: all docs at an agenda item → one LLM call
  Level 2 — Item rollup: parent items synthesise child summaries
  Level 3 — Meeting briefing: top-level items → meeting summary

Text extraction: .pdf (pymupdf/fitz), .docx (python-docx), .pptx (python-pptx)
LLM: Anthropic API via ANTHROPIC_API_KEY environment variable
"""
from __future__ import annotations

import hashlib
import json
import logging
import re
from pathlib import Path
from typing import Callable

import base64
import io

import fitz  # pymupdf
import yaml
from docx import Document as DocxDocument
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

import pipeline.db_new as db

logger = logging.getLogger(__name__)

SUMMARIZE_EXTENSIONS = {".pdf", ".docx", ".pptx"}

# Absolute path to repo root (parent of this file's parent)
_REPO_ROOT = Path(__file__).parent.parent
_PROMPTS_DIR = _REPO_ROOT / "prompts"

# Model IDs
HAIKU  = "claude-haiku-4-5-20251001"
SONNET = "claude-sonnet-4-6"
OPUS   = "claude-opus-4-6"

_DEFAULT_MODELS = {
    "document_model": HAIKU,
    "item_model":     HAIKU,
    "meeting_model":  HAIKU,
}

_DEFAULT_MAX_TOKENS = {
    "document_max_tokens": 32768,
    "item_max_tokens":     32768,
    "meeting_max_tokens":  32768,
}


# ---------------------------------------------------------------------------
# Anthropic client
# ---------------------------------------------------------------------------

def make_client():
    """Create an Anthropic client. Reads ANTHROPIC_API_KEY from env."""
    import anthropic
    return anthropic.Anthropic()


# ---------------------------------------------------------------------------
# Prompt and model config helpers
# ---------------------------------------------------------------------------

def _load_prompt(slug: str) -> str:
    path = _PROMPTS_DIR / f"{slug}.md"
    if path.exists():
        return path.read_text(encoding="utf-8")
    return ""


def _load_model_config() -> dict:
    path = _PROMPTS_DIR / "model_config.json"
    defaults = {**_DEFAULT_MODELS, **_DEFAULT_MAX_TOKENS}
    if path.exists():
        try:
            return {**defaults, **json.loads(path.read_text(encoding="utf-8"))}
        except Exception:
            pass
    return defaults


_VENUE_SLUG = {
    "ISO-NE": "isone",
    "NYISO":  "nyiso",
}


def _get_committee_prompts(
    committee_short: str,
    venue_short: str = "ISO-NE",
    briefing_style: str = "standard",
) -> tuple[str, str]:
    """
    Return (briefing_prompt, agenda_item_prompt) for the given committee
    short name (e.g. "MC", "RC", "NPC") and venue (e.g. "ISO-NE", "NYISO").

    briefing_style: "standard" or "detailed".  When "detailed", the lookup
    tries {venue}_{committee}_briefing_detailed_prompt first, falling back
    to the standard prompt if no detailed variant exists.

    Lookup order for each prompt type:
      1. {venue}_{committee}_{type}_prompt  (e.g. isone_mc_briefing_prompt)
      2. {committee}_{type}_prompt          (e.g. mc_briefing_prompt)
      3. isone_mc_{type}_prompt             (ultimate fallback)

    Prepends general_context_prompt if non-empty.
    """
    venue_slug = _VENUE_SLUG.get(venue_short, venue_short.lower().replace("-", ""))
    comm_slug = committee_short.lower()

    briefing = None
    if briefing_style == "detailed":
        briefing = (
            _load_prompt(f"{venue_slug}_{comm_slug}_briefing_detailed_prompt")
            or _load_prompt(f"{comm_slug}_briefing_detailed_prompt")
        )
    if not briefing:
        briefing = (
            _load_prompt(f"{venue_slug}_{comm_slug}_briefing_prompt")
            or _load_prompt(f"{comm_slug}_briefing_prompt")
            or _load_prompt("isone_mc_briefing_prompt")
        )
    agenda_item = (
        _load_prompt(f"{venue_slug}_{comm_slug}_agenda_item_prompt")
        or _load_prompt(f"{comm_slug}_agenda_item_prompt")
        or _load_prompt("isone_mc_agenda_item_prompt")
    )
    ctx = _load_prompt("general_context_prompt").strip()
    if ctx:
        briefing    = ctx + "\n\n" + briefing
        agenda_item = ctx + "\n\n" + agenda_item
    return briefing, agenda_item


# ---------------------------------------------------------------------------
# Text extraction
# ---------------------------------------------------------------------------

def extract_text_pdf(file_path: Path) -> str:
    doc = fitz.open(str(file_path))
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text()
        if text.strip():
            pages.append(f"[Page {i + 1}]\n{text.strip()}")
    doc.close()
    return "\n\n".join(pages)


def extract_text_docx(file_path: Path) -> str:
    doc = DocxDocument(str(file_path))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                parts.append(row_text)
    return "\n\n".join(parts)


def extract_text_pptx(file_path: Path) -> str:
    prs = Presentation(str(file_path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def extract_text(file_path: Path) -> str:
    ext = file_path.suffix.lower()
    if ext == ".pdf":
        return extract_text_pdf(file_path)
    if ext == ".docx":
        return extract_text_docx(file_path)
    if ext == ".pptx":
        return extract_text_pptx(file_path)
    raise ValueError(f"No extractor for file type: {ext}")


# ---------------------------------------------------------------------------
# Image extraction
# ---------------------------------------------------------------------------

def _load_image_config() -> dict:
    """Load image extraction settings from config.yaml."""
    cfg_path = _REPO_ROOT / "config.yaml"
    defaults = {
        "enabled": False,
        "min_size_px": 200,
        "max_per_item": 10,
        "extract_from": [".pdf", ".pptx"],
    }
    if cfg_path.exists():
        try:
            cfg = yaml.safe_load(cfg_path.read_text(encoding="utf-8"))
            return {**defaults, **cfg.get("summarization", {}).get("images", {})}
        except Exception:
            pass
    return defaults


def _img_to_png_b64(image_bytes: bytes) -> str:
    """Convert raw image bytes to a base64-encoded PNG string."""
    img = Image.open(io.BytesIO(image_bytes))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return base64.b64encode(buf.getvalue()).decode("ascii")


def extract_images_pptx(
    file_path: Path, min_px: int = 200,
) -> list[dict]:
    """
    Extract images from a PPTX file.
    Returns list of dicts: {page_or_slide, img_index, width, height,
                            image_bytes, media_type}.
    Deduplicates by SHA-256 of image bytes.
    """
    prs = Presentation(str(file_path))
    results = []
    seen_hashes: set[str] = set()

    for slide_num, slide in enumerate(prs.slides, 1):
        img_idx = 0
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                blob = shape.image.blob
                content_type = shape.image.content_type  # e.g. "image/png"
            except Exception:
                continue

            # Deduplicate
            img_hash = hashlib.sha256(blob).hexdigest()
            if img_hash in seen_hashes:
                continue
            seen_hashes.add(img_hash)

            # Size check
            try:
                img = Image.open(io.BytesIO(blob))
                w, h = img.size
            except Exception:
                continue
            if w < min_px and h < min_px:
                continue

            results.append({
                "page_or_slide": slide_num,
                "img_index": img_idx,
                "width": w,
                "height": h,
                "image_bytes": blob,
                "media_type": content_type or "image/png",
            })
            img_idx += 1

    return results


def extract_images_pdf(
    file_path: Path, min_px: int = 200,
) -> list[dict]:
    """
    Extract images from a PDF file using pymupdf.
    Returns same format as extract_images_pptx.
    """
    doc = fitz.open(str(file_path))
    results = []
    seen_hashes: set[str] = set()

    for page_num in range(len(doc)):
        page = doc[page_num]
        img_list = page.get_images(full=True)
        img_idx = 0
        for img_info in img_list:
            xref = img_info[0]
            try:
                pix = doc.extract_image(xref)
            except Exception:
                continue
            if not pix or not pix.get("image"):
                continue

            blob = pix["image"]
            ext = pix.get("ext", "png")
            w = pix.get("width", 0)
            h = pix.get("height", 0)

            # Deduplicate
            img_hash = hashlib.sha256(blob).hexdigest()
            if img_hash in seen_hashes:
                continue
            seen_hashes.add(img_hash)

            # Size check
            if w < min_px and h < min_px:
                continue

            media_type = {
                "png": "image/png",
                "jpeg": "image/jpeg",
                "jpg": "image/jpeg",
            }.get(ext, "image/png")

            results.append({
                "page_or_slide": page_num + 1,
                "img_index": img_idx,
                "width": w,
                "height": h,
                "image_bytes": blob,
                "media_type": media_type,
            })
            img_idx += 1

    doc.close()
    return results


def extract_images(file_path: Path, min_px: int = 200) -> list[dict]:
    """Dispatch image extraction by file type. Returns [] for unsupported types."""
    ext = file_path.suffix.lower()
    if ext == ".pptx":
        return extract_images_pptx(file_path, min_px=min_px)
    if ext == ".pdf":
        return extract_images_pdf(file_path, min_px=min_px)
    return []


def _extract_and_store_images(
    doc: dict,
    meeting_folder: Path | None = None,
    min_px: int = 200,
) -> list[dict]:
    """
    Extract images from a document, store on disk + DB, and return image dicts.
    Skips if images already exist in the DB for this document.
    """
    # Already extracted?
    existing_count = db.count_images_for_document(doc["id"])
    if existing_count > 0:
        return db.get_images_for_document(doc["id"], min_size=min_px)

    file_type = (doc.get("file_type") or "").lower()
    img_cfg = _load_image_config()
    if file_type not in img_cfg.get("extract_from", [".pdf", ".pptx"]):
        return []

    source_url = doc.get("source_url")
    if not source_url:
        return []

    try:
        import requests as _requests
        from pipeline.downloader import download_file_temp
        session = _requests.Session()
        with download_file_temp(
            url=source_url,
            filename=doc["filename"],
            referer_url=source_url,
            session=session,
        ) as tmp_path:
            if not tmp_path:
                return []
            raw_images = extract_images(Path(tmp_path), min_px=min_px)
    except Exception as exc:
        logger.warning("Failed to extract images for %s: %s", doc["filename"], exc)
        return []

    if not raw_images:
        return []

    # Build filesystem-safe document stem
    doc_stem = Path(doc["filename"]).stem.replace(" ", "_")

    # Prepare output directory
    images_dir = None
    if meeting_folder:
        images_dir = meeting_folder / "_images"
        images_dir.mkdir(parents=True, exist_ok=True)

    stored = []
    for img in raw_images:
        slug = f"slide{img['page_or_slide']}" if file_type == ".pptx" else f"page{img['page_or_slide']}"
        fname = f"{doc_stem}_{slug}_img{img['img_index']}.png"

        # Convert to PNG base64
        b64 = _img_to_png_b64(img["image_bytes"])

        # Save to disk
        rel_path = None
        if images_dir:
            out_path = images_dir / fname
            png_bytes = base64.b64decode(b64)
            out_path.write_bytes(png_bytes)
            rel_path = str(out_path.relative_to(meeting_folder.parent.parent))

        row = db.insert_document_image(
            document_id=doc["id"],
            filename=fname,
            page_or_slide=img["page_or_slide"],
            img_index=img["img_index"],
            width=img["width"],
            height=img["height"],
            file_path=rel_path,
            image_b64=b64,
        )
        stored.append(row)

    logger.info("Extracted %d image(s) from %s", len(stored), doc["filename"])
    return stored


def file_sha256(file_path: Path) -> str:
    h = hashlib.sha256()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def _get_text_for_doc(doc: dict) -> str:
    """
    Return extracted plain text for a document row.
    Uses pre-stored raw_content if available; otherwise downloads and
    extracts the file, caching the result in raw_content for future use.
    Returns empty string if the document cannot be extracted.
    """
    # Use cached text if available
    if doc.get("raw_content"):
        return doc["raw_content"]

    # Skip CEII-protected or unsupported types
    file_type = (doc.get("file_type") or "").lower()
    if file_type not in SUMMARIZE_EXTENSIONS:
        logger.debug("Skipping %s — unsupported type %s", doc["filename"], file_type)
        return ""

    source_url = doc.get("source_url")
    if not source_url:
        logger.warning("No source_url for document %s — cannot extract text", doc["filename"])
        return ""

    try:
        import requests as _requests
        from pipeline.downloader import download_file_temp
        session = _requests.Session()
        with download_file_temp(
            url=source_url,
            filename=doc["filename"],
            referer_url=source_url,
            session=session,
        ) as tmp_path:
            if not tmp_path:
                logger.warning("Download returned None for %s", doc["filename"])
                return ""
            text = extract_text(Path(tmp_path))
            # Cache the result in DB
            if text:
                db.set_document_raw_content(doc["id"], text)
            return text
    except Exception as exc:
        logger.warning("Failed to extract text for %s: %s", doc["filename"], exc)
        return ""


# ---------------------------------------------------------------------------
# Agenda metadata helpers
# ---------------------------------------------------------------------------

def _item_metadata_block(item: dict) -> str:
    """
    Build a compact metadata block for an agenda item using its stored fields.
    Returns an empty string if no metadata is present.
    """
    lines = []
    label = item.get("item_id") or ""
    title = item.get("title") or ""
    lines.append(f"**Agenda item:** {label}  {title}".strip())

    presenter = item.get("presenter")
    org       = item.get("org")
    if presenter:
        lines.append(f"**Presenter:** {presenter}" + (f" ({org})" if org else ""))
    elif org:
        lines.append(f"**Organisation:** {org}")

    vote = item.get("vote_status")
    if vote:
        lines.append(f"**Vote/action status:** {vote}")

    time_slot = item.get("time_slot")
    if time_slot:
        lines.append(f"**Scheduled time:** {time_slot}")

    wmpp = item.get("wmpp_id")
    if wmpp:
        lines.append(f"**WMPP/workstream:** {wmpp}")

    notes = item.get("notes")
    if notes:
        lines.append(f"**Notes:** {notes}")

    if len(lines) <= 1:
        return ""  # Only the title line — no useful metadata
    return "\n".join(lines)


def _meeting_structure_block(all_items: list[dict]) -> str:
    """
    Build a hierarchical agenda outline for all items in the meeting.
    Includes presenter, org, vote status, and time slot where available.
    Used as context in the Level 3 briefing prompt.
    """
    lines = ["**Meeting Agenda Structure:**\n"]
    for item in all_items:
        depth     = item.get("depth", 0)
        indent    = "  " * depth
        label     = item.get("item_id") or ""
        title     = item.get("title") or ""
        presenter = item.get("presenter") or ""
        org       = item.get("org") or ""
        vote      = item.get("vote_status") or ""
        time_slot = item.get("time_slot") or ""

        meta_parts = []
        if presenter:
            meta_parts.append(presenter + (f" / {org}" if org else ""))
        if vote:
            meta_parts.append(vote)
        if time_slot:
            meta_parts.append(time_slot)

        meta_str = f"  ·  {',  '.join(meta_parts)}" if meta_parts else ""
        lines.append(f"{indent}- **{label}**  {title}{meta_str}".rstrip())

    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Output cleaning
# ---------------------------------------------------------------------------

# Patterns whose titles should be excluded from the Level 3 briefing
_SKIP_BRIEFING_RE = re.compile(
    r"\b(meeting\s+minutes?|draft\s+minutes?|minutes?\s+approval|"
    r"approval\s+of\s+minutes?|attendance|roll\s+call|"
    r"chair['\u2019]?s?\s+(opening\s+)?remarks?|opening\s+remarks?|"
    r"administrative)\b",
    re.IGNORECASE,
)


def _should_skip_briefing(item: dict) -> bool:
    """Return True if an agenda item should be excluded from the Level 3 briefing."""
    title = (item.get("title") or "")
    return bool(_SKIP_BRIEFING_RE.search(title))


def _clean_output(text: str) -> str:
    """
    Post-process LLM output for safe Markdown rendering.

    - Escapes bare $ signs (currency/math) so Streamlit does not interpret
      them as LaTeX delimiters. Handles both "$32.3M" and "$...$" pairs.
    - Leaves already-escaped \\$ untouched.
    """
    # Escape any $ that is not already preceded by a backslash
    return re.sub(r"(?<!\\)\$", r"\\$", text)


# ---------------------------------------------------------------------------
# LLM call helper
# ---------------------------------------------------------------------------

def _call_llm(client, model: str, prompt: str, max_tokens: int = 4096,
              max_retries: int = 3, label: str = "") -> str:
    """Streaming LLM call with retry on rate-limit and truncation handling."""
    import time as _time
    current_max = max_tokens
    for attempt in range(max_retries):
        try:
            with client.messages.stream(
                model=model,
                max_tokens=current_max,
                messages=[{"role": "user", "content": prompt}],
            ) as stream:
                msg = stream.get_final_message()
            text = _clean_output(msg.content[0].text.strip())
            if msg.stop_reason == "max_tokens":
                retry_max = min(current_max * 2, 65536)
                if retry_max > current_max:
                    logger.warning(
                        "Output truncated at %d tokens%s — retrying with %d",
                        current_max, f" ({label})" if label else "", retry_max,
                    )
                    current_max = retry_max
                    continue
                logger.warning(
                    "Output still truncated at %d tokens%s — returning partial result",
                    current_max, f" ({label})" if label else "",
                )
            return text
        except Exception as exc:
            if "rate_limit" in str(exc).lower() and attempt < max_retries - 1:
                wait = 30 * (attempt + 1)
                logger.warning("Rate limited (attempt %d/%d), waiting %ds...",
                               attempt + 1, max_retries, wait)
                _time.sleep(wait)
            else:
                raise


def _call_llm_multimodal(
    client, model: str, text_prompt: str,
    images: list[dict],
    max_tokens: int = 4096,
    max_retries: int = 3,
    max_images: int = 10,
    label: str = "",
) -> str:
    """
    Multimodal LLM call with text + images.
    Falls back to text-only if no images provided.
    Images are sorted by area descending (largest first) and capped at max_images.
    """
    if not images:
        return _call_llm(client, model, text_prompt, max_tokens=max_tokens,
                         max_retries=max_retries, label=label)

    # Sort by area descending — larger images are more likely substantive
    sorted_imgs = sorted(
        images,
        key=lambda x: (x.get("width") or 0) * (x.get("height") or 0),
        reverse=True,
    )[:max_images]

    # Build multimodal content blocks
    content: list[dict] = [{"type": "text", "text": text_prompt}]
    for img in sorted_imgs:
        b64 = img.get("image_b64") or ""
        if not b64:
            continue
        media_type = img.get("media_type", "image/png")
        # Normalise media type for the API
        if media_type not in ("image/png", "image/jpeg", "image/gif", "image/webp"):
            media_type = "image/png"
        doc_name = img.get("doc_filename") or img.get("filename") or "image"
        page_label = f"slide {img['page_or_slide']}" if "slide" in (img.get("filename") or "") else f"page {img['page_or_slide']}"
        label_idx = img.get("_label_idx", "?")
        content.append({
            "type": "image",
            "source": {
                "type": "base64",
                "media_type": media_type,
                "data": b64,
            },
        })
        content.append({
            "type": "text",
            "text": f"[Image {label_idx} from {doc_name}, {page_label}]",
        })

    import time as _time
    current_max = max_tokens
    for attempt in range(max_retries):
        try:
            with client.messages.stream(
                model=model,
                max_tokens=current_max,
                messages=[{"role": "user", "content": content}],
            ) as stream:
                msg = stream.get_final_message()
            text = _clean_output(msg.content[0].text.strip())
            if msg.stop_reason == "max_tokens":
                retry_max = min(current_max * 2, 65536)
                if retry_max > current_max:
                    logger.warning(
                        "Output truncated at %d tokens%s — retrying with %d",
                        current_max, f" ({label})" if label else "", retry_max,
                    )
                    current_max = retry_max
                    continue
            return text
        except Exception as exc:
            if "rate_limit" in str(exc).lower() and attempt < max_retries - 1:
                wait = 30 * (attempt + 1)
                logger.warning("Rate limited (attempt %d/%d), waiting %ds...",
                               attempt + 1, max_retries, wait)
                _time.sleep(wait)
            else:
                raise


# ---------------------------------------------------------------------------
# Level 1 — Document-group summary (per agenda item)
# ---------------------------------------------------------------------------

def _run_item_doc_summary(
    item: dict,
    client,
    model: str,
    doc_summary_prompt: str,
    max_tokens: int = 4096,
    extract_images: bool = False,
    meeting_folder: Path | None = None,
) -> bool:
    """
    Summarise all documents assigned to `item` in one LLM call.
    Stores result as a new summary_version (entity_type='agenda_item').
    Returns True if a summary was created, False if no usable docs found.
    """
    docs = db.get_documents_for_item(item["id"])
    usable = [
        d for d in docs
        if not d.get("ceii_skipped")
        and not d.get("ignored")
        and (d.get("file_type") or "").lower() in SUMMARIZE_EXTENSIONS
    ]

    if not usable:
        logger.debug("Item %s (%s): no usable documents — skipping Level 1",
                     item.get("item_id"), item["title"])
        return False

    # Extract text for each doc (and images if enabled)
    text_parts = []
    all_images: list[dict] = []
    img_cfg = _load_image_config()

    for doc in usable:
        text = _get_text_for_doc(doc)
        if text:
            text_parts.append(f"### [{doc['filename']}]\n\n{text}")

        # Image extraction (opt-in)
        if extract_images:
            try:
                doc_images = _extract_and_store_images(
                    doc,
                    meeting_folder=meeting_folder,
                    min_px=img_cfg.get("min_size_px", 200),
                )
                all_images.extend(doc_images)
            except Exception as exc:
                logger.warning("Image extraction failed for %s: %s", doc["filename"], exc)

    if not text_parts:
        logger.warning("Item %s (%s): all docs returned empty text",
                       item.get("item_id"), item["title"])
        return False

    combined_text      = "\n\n---\n\n".join(text_parts)
    combined_filenames = ", ".join(d["filename"] for d in usable if d["filename"])
    item_label         = item.get("item_id") or item["title"]

    # Build metadata context block and prepend to the document text
    meta_block = _item_metadata_block(item)
    if meta_block:
        augmented_text = f"{meta_block}\n\n---\n\n{combined_text}"
    else:
        augmented_text = combined_text

    # Fill the doc_summary_prompt template
    # Supports {filename} and {text} placeholders
    try:
        prompt = doc_summary_prompt.format(
            filename=combined_filenames,
            text=augmented_text,
        )
    except KeyError:
        prompt = doc_summary_prompt + f"\n\n---\n\nDocument(s): {combined_filenames}\n\n{augmented_text}"

    # Append visual content instruction when images are available
    if extract_images and all_images:
        # Number images so the model can reference them by index
        for idx, img in enumerate(all_images):
            img["_label_idx"] = idx + 1
        prompt += (
            "\n\n---\n\n## Visual Content\n"
            "Images from the document(s) are provided below, numbered sequentially. "
            "Use the visual data to enrich your summary — weave chart insights into "
            "the narrative sections above.\n\n"
            "Select **0 to 2** 'killer images' — charts, diagrams, or tables that "
            "provide irreplaceable visual context a reader absolutely needs to see. "
            "Place each KEEP_IMAGE directive **inline** in your summary, on its own "
            "line, right after the paragraph where the image is most relevant:\n"
            "KEEP_IMAGE <N>: <one-sentence caption>\n\n"
            "where <N> is the image number shown in [Image N from ...] labels. "
            "Do NOT group them at the end. If no image is worth including, simply "
            "omit any KEEP_IMAGE lines."
        )

    logger.info("Level 1 — item %s: summarising %d doc(s)%s with %s",
                item_label, len(text_parts),
                f" + {len(all_images)} image(s)" if all_images else "",
                model)

    max_imgs = img_cfg.get("max_per_item", 10)
    if extract_images and all_images:
        detailed = _call_llm_multimodal(
            client, model, prompt, all_images,
            max_tokens=max_tokens, max_images=max_imgs,
            label=f"L1 item {item_label}",
        )
    else:
        detailed = _call_llm(client, model, prompt, max_tokens=max_tokens,
                             label=f"L1 item {item_label}")

    # Post-process: replace KEEP_IMAGE lines inline with image embeds
    if extract_images and all_images and detailed:
        detailed = _replace_keep_images_inline(detailed, all_images)

    db.create_summary_version(
        entity_type="agenda_item",
        entity_id=item["id"],
        one_line=None,
        detailed=detailed,
        model_id=model,
        is_manual=False,
        status="draft",
        created_by="system",
    )
    return True


def _extract_kept_images(llm_output: str, all_images: list[dict]) -> list[dict]:
    """
    Parse KEEP_IMAGE lines from LLM output.
    Returns up to 2 image dicts from all_images, annotated with _caption.
    """
    kept = []
    for m in re.finditer(r"KEEP_IMAGE\s+(\d+)\s*:\s*(.+)", llm_output):
        idx = int(m.group(1))
        caption = m.group(2).strip()
        # Images were numbered 1-based via _label_idx
        match = next((img for img in all_images if img.get("_label_idx") == idx), None)
        if match:
            match["_caption"] = caption
            kept.append(match)
        if len(kept) >= 2:
            break
    return kept


# ---------------------------------------------------------------------------
# Image reference helpers (used by L2 rollup and L3 briefing)
# ---------------------------------------------------------------------------

def _collect_image_refs(text: str) -> list[int]:
    """Extract image_id integers from <!-- image_id:N --> comments in summary text."""
    return [int(m) for m in re.findall(r"<!-- image_id:(\d+) -->", text)]


def _fetch_images_for_refs(image_ids: list[int]) -> list[dict]:
    """
    Fetch image records from DB for a list of image_ids.
    Decodes base64 back to raw bytes for multimodal LLM calls.
    Returns list of dicts compatible with _call_llm_multimodal().
    """
    if not image_ids:
        return []
    rows = db.get_images_by_ids(image_ids)
    images = []
    for row in rows:
        if not row.get("image_b64"):
            continue
        raw_bytes = base64.b64decode(row["image_b64"])
        images.append({
            "id": row["id"],
            "image_bytes": raw_bytes,
            "page_or_slide": row.get("page_or_slide", 0),
            "width": row.get("width", 0),
            "height": row.get("height", 0),
            "filename": row.get("filename", ""),
            "description": row.get("description", ""),
        })
    return images


_IMAGE_ROLLUP_INSTRUCTION = (
    "\n\n---\n\n## Visual Content\n"
    "Images below are key figures selected from the underlying document summaries. "
    "They have already been filtered to the most substantive charts and diagrams.\n\n"
    "You may select **0 to 2** of these images to carry forward — only the "
    "'killer image' that provides irreplaceable visual context. "
    "Place each KEEP_IMAGE directive **inline** in your summary, on its own line, "
    "right after the paragraph where the image is most relevant:\n"
    "KEEP_IMAGE <N>: <one-sentence caption>\n\n"
    "where <N> is the image number shown in [Image N ...] labels. "
    "Do NOT group them at the end. "
    "If no image adds value beyond what the text conveys, omit any KEEP_IMAGE lines."
)

_IMAGE_BRIEFING_INSTRUCTION = (
    "\n\n---\n\n## Visual Content\n"
    "Images below are key figures selected from the agenda item summaries. "
    "They have already been filtered to the most substantive charts and diagrams.\n\n"
    "You may select **0 to 2** images for the entire briefing — only the "
    "'killer chart' that anchors understanding of a key development. "
    "Place each KEEP_IMAGE directive **inline** in the relevant agenda item "
    "section, on its own line, right after the paragraph it illustrates:\n"
    "KEEP_IMAGE <N>: <one-sentence caption>\n\n"
    "where <N> is the image number shown in [Image N ...] labels. "
    "Do NOT group them at the end or in a separate figures section. "
    "If no image adds value beyond what the text conveys, omit any KEEP_IMAGE lines."
)


def _replace_keep_images_inline(detailed: str, all_images: list[dict]) -> str:
    """
    Replace KEEP_IMAGE <N>: <caption> lines in-place with inline image embeds.
    Each KEEP_IMAGE line becomes:
        **Figure:** <caption>
        <!-- image_id:<id> -->
    Preserves the image's position in the narrative. Up to 2 kept.
    """
    if not all_images or not detailed:
        return detailed

    kept_count = 0

    def _replace_match(m: re.Match) -> str:
        nonlocal kept_count
        if kept_count >= 2:
            return ""  # strip extras beyond 2
        idx = int(m.group(1))
        caption = m.group(2).strip()
        # Find the matching image by _label_idx first, then fall back to
        # matching by DB id (the LLM sometimes uses the actual DB id from
        # <!-- image_id:N --> comments in the source summaries).
        match = next((img for img in all_images if img.get("_label_idx") == idx), None)
        if not match:
            match = next((img for img in all_images if img.get("id") == idx), None)
        if not match:
            # Image may be from a child item not in the pre-collected set;
            # look it up directly in the DB.
            fetched = _fetch_images_for_refs([idx])
            match = fetched[0] if fetched else None
        if not match:
            return ""
        kept_count += 1
        img_id = match.get("id")
        # Store caption in DB
        if img_id:
            try:
                db.set_image_description(img_id, caption)
            except Exception:
                pass
        return f"\n**Figure:** {caption}\n\n<!-- image_id:{img_id} -->\n"

    detailed = re.sub(
        r"^KEEP_IMAGE\s+<?(\d+)>?\s*:\s*(.+)$",
        _replace_match,
        detailed,
        flags=re.MULTILINE,
    )
    # Clean up any leftover KEEP_IMAGE lines (beyond the 2 kept)
    detailed = re.sub(r"^KEEP_IMAGE\s+<?\d+>?:.*$", "", detailed, flags=re.MULTILINE)
    # Remove empty Key Figures / Visual Content headings the LLM may have added
    detailed = re.sub(r"##\s*(?:Key Figures|Visual Content)\s*\n*$", "", detailed)
    # Remove "No key figures" lines
    detailed = re.sub(r"^No key figures\.?\s*$", "", detailed, flags=re.MULTILINE)
    return detailed.strip()


# ---------------------------------------------------------------------------
# Public wrappers — used by pipeline/deep_dive.py
# ---------------------------------------------------------------------------

def get_text_for_doc(doc: dict) -> str:
    """Public wrapper for _get_text_for_doc."""
    return _get_text_for_doc(doc)


def extract_and_store_images(
    doc: dict,
    meeting_folder: Path | None = None,
    min_px: int = 200,
) -> list[dict]:
    """Public wrapper for _extract_and_store_images."""
    return _extract_and_store_images(doc, meeting_folder=meeting_folder, min_px=min_px)


def call_llm_multimodal(
    client, model: str, text_prompt: str,
    images: list[dict],
    max_tokens: int = 4096,
    max_retries: int = 3,
    max_images: int = 10,
    label: str = "",
) -> str:
    """Public wrapper for _call_llm_multimodal."""
    return _call_llm_multimodal(
        client, model, text_prompt, images,
        max_tokens=max_tokens, max_retries=max_retries,
        max_images=max_images, label=label,
    )


def replace_keep_images_inline(detailed: str, all_images: list[dict],
                                max_keep: int = 2) -> str:
    """
    Public wrapper for _replace_keep_images_inline.
    Supports a configurable max_keep (default 2 for backward compatibility).
    """
    if max_keep == 2:
        return _replace_keep_images_inline(detailed, all_images)
    # For custom max_keep, do inline replacement with different cap
    if not all_images or not detailed:
        return detailed
    kept_count = 0

    def _replace_match(m: re.Match) -> str:
        nonlocal kept_count
        if kept_count >= max_keep:
            return ""
        idx = int(m.group(1))
        caption = m.group(2).strip()
        match = next((img for img in all_images if img.get("_label_idx") == idx), None)
        if not match:
            match = next((img for img in all_images if img.get("id") == idx), None)
        if not match:
            fetched = _fetch_images_for_refs([idx])
            match = fetched[0] if fetched else None
        if not match:
            return ""
        kept_count += 1
        img_id = match.get("id")
        if img_id:
            try:
                db.set_image_description(img_id, caption)
            except Exception:
                pass
        return f"\n**Figure:** {caption}\n\n<!-- image_id:{img_id} -->\n"

    detailed = re.sub(
        r"^KEEP_IMAGE\s+<?(\d+)>?\s*:\s*(.+)$",
        _replace_match, detailed, flags=re.MULTILINE,
    )
    detailed = re.sub(r"^KEEP_IMAGE\s+<?\d+>?:.*$", "", detailed, flags=re.MULTILINE)
    detailed = re.sub(r"##\s*(?:Key Figures|Visual Content)\s*\n*$", "", detailed)
    detailed = re.sub(r"^No key figures\.?\s*$", "", detailed, flags=re.MULTILINE)
    return detailed.strip()


def clean_output(text: str) -> str:
    """Public wrapper for _clean_output."""
    return _clean_output(text)


def load_prompt(slug: str) -> str:
    """Public wrapper for _load_prompt."""
    return _load_prompt(slug)


def load_model_config() -> dict:
    """Public wrapper for _load_model_config."""
    return _load_model_config()


# ---------------------------------------------------------------------------
# Level 2 — Item rollup (parent items with child summaries)
# ---------------------------------------------------------------------------

def _run_item_rollup(
    item: dict,
    child_summaries: list[dict],
    client,
    model: str,
    agenda_item_prompt: str,
    max_tokens: int = 4096,
) -> bool:
    """
    Synthesise child summaries into a summary for `item`.
    Stores result as a new summary_version (entity_type='agenda_item').
    Returns True if a summary was created.
    """
    parts = []
    for child, summ in child_summaries:
        child_label = child.get("item_id") or child["title"]
        detailed    = summ.get("detailed") or ""
        if detailed:
            parts.append(f"**{child_label} — {child['title']}**\n\n{detailed}")

    if not parts:
        return False

    doc_summaries_block = "\n\n---\n\n".join(parts)
    item_label = item.get("item_id") or item["title"]

    # Prepend item metadata so the model knows presenter/org/vote context
    meta_block = _item_metadata_block(item)
    if meta_block:
        doc_summaries_block = f"{meta_block}\n\n---\n\n{doc_summaries_block}"

    try:
        prompt = agenda_item_prompt.format(
            item_id=item_label,
            title=item["title"],
            doc_summaries=doc_summaries_block,
        )
    except KeyError:
        prompt = (
            agenda_item_prompt
            + f"\n\nAgenda item: {item_label} — {item['title']}\n\n"
            + doc_summaries_block
        )

    # Collect images referenced in child summaries
    all_image_ids: list[int] = []
    for child, summ in child_summaries:
        detailed_text = summ.get("detailed") or ""
        all_image_ids.extend(_collect_image_refs(detailed_text))
    all_images = _fetch_images_for_refs(all_image_ids)

    if all_images:
        # Number images for multimodal labeling
        for idx, img in enumerate(all_images):
            img["_label_idx"] = idx + 1
        prompt += _IMAGE_ROLLUP_INSTRUCTION

    logger.info("Level 2 — item %s: rolling up %d child summaries%s with %s",
                item_label, len(parts),
                f" + {len(all_images)} image(s)" if all_images else "",
                model)

    if all_images:
        detailed = _call_llm_multimodal(
            client, model, prompt, all_images,
            max_tokens=max_tokens, max_images=8,
            label=f"L2 item {item_label}",
        )
    else:
        detailed = _call_llm(client, model, prompt, max_tokens=max_tokens,
                             label=f"L2 item {item_label}")

    detailed = _replace_keep_images_inline(detailed, all_images)

    db.create_summary_version(
        entity_type="agenda_item",
        entity_id=item["id"],
        one_line=None,
        detailed=detailed,
        model_id=model,
        is_manual=False,
        status="draft",
        created_by="system",
    )
    return True


# ---------------------------------------------------------------------------
# Level 3 — Meeting briefing
# ---------------------------------------------------------------------------

def _run_meeting_briefing(
    meeting_id: int,
    top_level_items: list[dict],
    all_items: list[dict],
    client,
    model: str,
    briefing_prompt: str,
    max_tokens: int = 8192,
) -> bool:
    """
    Synthesise top-level item summaries into the meeting briefing.
    Prepends the full agenda structure (with presenter/org/vote metadata)
    as context before the item summaries.
    Stores result as entity_type='meeting'.
    Returns True if briefing was created.
    """
    parts = []
    for item in top_level_items:
        # Skip administrative items (minutes, opening remarks, etc.)
        if _should_skip_briefing(item):
            logger.info("Level 3 — skipping administrative item: %s",
                        item.get("title"))
            continue
        summ = db.get_current_summary("agenda_item", item["id"])
        if summ and summ.get("detailed"):
            label = item.get("item_id") or item["title"]
            # Include per-item metadata header in each block
            meta = _item_metadata_block(item)
            header = f"## Item {label}: {item['title']}"
            if meta:
                header += f"\n{meta}"
            parts.append(f"{header}\n\n{summ['detailed']}")

    if not parts:
        logger.warning("Level 3 — meeting %d: no item summaries available for briefing",
                       meeting_id)
        return False

    # Prepend the full agenda structure so the model understands
    # overall meeting themes, flow, presenters, and action items
    structure_block = _meeting_structure_block(all_items)
    items_block     = "\n\n---\n\n".join(parts)
    context_block   = f"{structure_block}\n\n---\n\n{items_block}"

    if "[AGENDA ITEMS]" in briefing_prompt:
        prompt = briefing_prompt.replace("[AGENDA ITEMS]", context_block)
    else:
        prompt = briefing_prompt + "\n\n" + context_block

    # Collect images referenced across all item summaries
    all_image_ids: list[int] = []
    for part_text in parts:
        all_image_ids.extend(_collect_image_refs(part_text))
    all_images = _fetch_images_for_refs(all_image_ids)

    if all_images:
        for idx, img in enumerate(all_images):
            img["_label_idx"] = idx + 1
        prompt += _IMAGE_BRIEFING_INSTRUCTION

    logger.info("Level 3 — meeting %d: generating briefing from %d items%s with %s",
                meeting_id, len(parts),
                f" + {len(all_images)} image(s)" if all_images else "",
                model)

    if all_images:
        detailed = _call_llm_multimodal(
            client, model, prompt, all_images,
            max_tokens=max_tokens, max_images=8,
            label=f"L3 meeting {meeting_id}",
        )
    else:
        detailed = _call_llm(client, model, prompt, max_tokens=max_tokens,
                             label=f"L3 meeting {meeting_id}")

    detailed = _replace_keep_images_inline(detailed, all_images)

    db.create_summary_version(
        entity_type="meeting",
        entity_id=meeting_id,
        one_line=None,
        detailed=detailed,
        model_id=model,
        is_manual=False,
        status="draft",
        created_by="system",
    )
    return True


# ---------------------------------------------------------------------------
# Orchestration
# ---------------------------------------------------------------------------

def run_meeting_summarization(
    meeting_id: int,
    client,
    committee_short: str = "MC",
    venue_short: str = "ISO-NE",
    progress_fn: Callable[[str], None] | None = None,
    force_rerun: bool = False,
    start_level: int = 1,
    extract_images: bool | None = None,
    briefing_style: str = "standard",
    item_ids: set[int] | None = None,
) -> dict:
    """
    Run the three-level summarization pipeline for a meeting.

    start_level controls where the pipeline begins:
        1 = full run (doc summaries → rollups → briefing)
        2 = skip Level 1 doc summaries; only rollups + briefing
        3 = briefing only

    extract_images: None = use config.yaml default, True/False = override.

    item_ids: if provided, only process these agenda item IDs at Level 1,
    and only re-roll-up their parents at Level 2. Level 3 briefing is
    always re-run when item_ids is set (since content changed).
    When None, all items are processed (default behavior).

    Returns a dict with counts:
        {"level1": int, "level2": int, "level3": bool, "errors": list[str]}

    progress_fn(message) is called with status strings for UI display.
    """

    def _progress(msg: str) -> None:
        logger.info(msg)
        if progress_fn:
            progress_fn(msg)

    # Resolve image extraction flag
    img_cfg = _load_image_config()
    do_images = extract_images if extract_images is not None else img_cfg.get("enabled", False)

    # Resolve meeting folder for image storage
    meeting_folder: Path | None = None
    if do_images:
        meeting = db.get_meeting(meeting_id)
        if meeting:
            venue = meeting.get("venue_short", "ISO-NE")
            cfg_path = _REPO_ROOT / "config.yaml"
            try:
                full_cfg = yaml.safe_load(cfg_path.read_text(encoding="utf-8"))
            except Exception:
                full_cfg = {}
            if venue == "NYISO":
                storage_root = Path(full_cfg.get("nyiso_storage_root", "./nyiso-materials"))
            else:
                storage_root = Path(full_cfg.get("storage_root", "./nepool-materials"))
            # Find existing meeting folder by scanning the committee directory
            type_short = meeting.get("type_short", "")
            type_name = meeting.get("type_name", "")
            comm_dir = _REPO_ROOT / storage_root / type_name
            if comm_dir.is_dir():
                for d in comm_dir.iterdir():
                    if d.is_dir() and type_short in d.name:
                        # Check if the meeting date is in the folder name
                        date_str = str(meeting["meeting_date"])
                        if date_str in d.name or meeting["meeting_date"].strftime("%m%d%y") in d.name:
                            meeting_folder = d
                            break
        if do_images:
            _progress(f"Image extraction enabled{' (folder: ' + meeting_folder.name + ')' if meeting_folder else ''}.")

    models      = _load_model_config()
    doc_model   = models["document_model"]
    item_model  = models["item_model"]
    mtg_model   = models["meeting_model"]
    doc_max_tokens  = models["document_max_tokens"]
    item_max_tokens = models["item_max_tokens"]
    mtg_max_tokens  = models["meeting_max_tokens"]

    briefing_prompt, agenda_item_prompt = _get_committee_prompts(committee_short, venue_short, briefing_style=briefing_style)
    doc_summary_prompt = _load_prompt("doc_summary_prompt")

    if not doc_summary_prompt:
        doc_summary_prompt = (
            "Summarise the following document(s) for an energy market analyst.\n\n"
            "Document(s): {filename}\n\n{text}"
        )

    all_items = db.get_agenda_items(meeting_id)
    if not all_items:
        _progress("No agenda items found — nothing to summarise.")
        return {"level1": 0, "level2": 0, "level3": False, "errors": []}

    # Build parent-child map — two mechanisms, merged together:
    #   1. Explicit parent_id FK (set during ingest)
    #   2. Prefix matching on item_id strings (fallback for "silent" parents
    #      whose children may have parent_id=NULL e.g. 2.1 → {2.1.a, 2.1.b, 2.1.c})
    item_by_id:      dict[int, dict]  = {it["id"]: it for it in all_items}
    item_by_item_id: dict[str, dict]  = {
        it["item_id"]: it for it in all_items if it.get("item_id")
    }
    children_of: dict[int, list[dict]] = {}

    for it in all_items:
        # Method 1 — explicit FK
        if it.get("parent_id") and it["parent_id"] in item_by_id:
            children_of.setdefault(it["parent_id"], [])
            if it not in children_of[it["parent_id"]]:
                children_of[it["parent_id"]].append(it)
        # Method 2 — prefix fallback (only when parent_id not set)
        # Walk UP the dot-separated prefix chain until we find an existing row.
        # e.g. "2.1.a" tries "2.1" first; if no "2.1" row exists, tries "2".
        elif it.get("item_id") and "." in it["item_id"]:
            parts = it["item_id"].split(".")
            parent_item = None
            for i in range(len(parts) - 1, 0, -1):
                candidate = ".".join(parts[:i])
                parent_item = item_by_item_id.get(candidate)
                if parent_item:
                    break
            if parent_item:
                children_of.setdefault(parent_item["id"], [])
                if it not in children_of[parent_item["id"]]:
                    children_of[parent_item["id"]].append(it)

    errors: list[str] = []
    l1_count = 0
    l2_count = 0

    # ── Level 1: leaf items (items with no children) ─────────────────────────
    if start_level > 1:
        _progress("Level 1: skipped (start_level > 1).")
    else:
        if item_ids:
            _progress(f"Level 1: summarising {len(item_ids)} affected agenda item(s)…")
        else:
            _progress("Level 1: summarising document groups for leaf agenda items…")
        for item in all_items:
            # If item_ids filter is set, skip items not in the set
            if item_ids is not None and item["id"] not in item_ids:
                continue

            docs = db.get_documents_for_item(item["id"])
            has_docs = any(
                not d.get("ceii_skipped") and not d.get("ignored")
                and (d.get("file_type") or "").lower() in SUMMARIZE_EXTENSIONS
                for d in docs
            )
            if not has_docs:
                continue

            # Skip if a non-stub summary already exists (avoids redundant LLM calls on re-run)
            existing = db.get_current_summary("agenda_item", item["id"])
            if (not force_rerun
                    and existing
                    and existing.get("status") not in ("stub", None)
                    and existing.get("detailed")):
                label = item.get("item_id") or item["title"]
                _progress(f"  Skipping item {label} — summary already exists (v{existing['version']})")
                l1_count += 1
                continue

            label = item.get("item_id") or item["title"]
            _progress(f"  Summarising documents for item {label}…")
            try:
                created = _run_item_doc_summary(
                    item, client, doc_model, doc_summary_prompt,
                    max_tokens=doc_max_tokens,
                    extract_images=do_images,
                    meeting_folder=meeting_folder,
                )
                if created:
                    l1_count += 1
            except Exception as exc:
                msg = f"Level 1 error on item {label}: {exc}"
                logger.error(msg)
                errors.append(msg)

    # ── Level 2: parent items (items that have children with summaries) ───────
    _progress("Level 2: rolling up parent items…")

    # Process parents in reverse depth order (deepest parents first)
    # so that when we roll up a grandparent, its children's rollups are ready
    parent_items = [it for it in all_items if children_of.get(it["id"])]
    parent_items_sorted = sorted(parent_items, key=lambda x: -x.get("depth", 0))

    # When item_ids filter is set, only re-roll-up parents whose children
    # include affected items
    affected_parent_ids: set[int] | None = None
    if item_ids is not None:
        affected_parent_ids = set()
        for parent in parent_items_sorted:
            children = children_of.get(parent["id"], [])
            if any(c["id"] in item_ids for c in children):
                affected_parent_ids.add(parent["id"])

    for item in parent_items_sorted:
        # Skip parents not affected by the item_ids filter
        if affected_parent_ids is not None and item["id"] not in affected_parent_ids:
            continue

        children = children_of.get(item["id"], [])
        child_summaries = []
        for child in children:
            summ = db.get_current_summary("agenda_item", child["id"])
            if summ and summ.get("detailed"):
                child_summaries.append((child, summ))

        if not child_summaries:
            continue

        label = item.get("item_id") or item["title"]

        # Skip if this parent already has a rollup summary
        existing = db.get_current_summary("agenda_item", item["id"])
        if (not force_rerun
                and existing
                and existing.get("status") not in ("stub", None)
                and existing.get("detailed")):
            _progress(f"  Skipping rollup for {label} — already exists (v{existing['version']})")
            l2_count += 1
            continue

        _progress(f"  Rolling up item {label} from {len(child_summaries)} child(ren)…")
        try:
            created = _run_item_rollup(
                item, child_summaries, client, item_model, agenda_item_prompt,
                max_tokens=item_max_tokens,
            )
            if created:
                l2_count += 1
        except Exception as exc:
            msg = f"Level 2 error on item {label}: {exc}"
            logger.error(msg)
            errors.append(msg)

    # ── Level 3: meeting briefing ─────────────────────────────────────────────
    _progress("Level 3: generating meeting briefing…")
    # When item_ids is set, always re-run briefing (affected content changed)
    l3_force = force_rerun or (item_ids is not None)
    existing_mtg = db.get_current_summary("meeting", meeting_id)
    if (not l3_force
            and existing_mtg
            and existing_mtg.get("status") not in ("stub", None)
            and existing_mtg.get("detailed")):
        _progress(f"  Skipping meeting briefing — already exists (v{existing_mtg['version']})")
        db.set_meeting_status(meeting_id, "complete")
        _progress("Done.")
        return {"level1": l1_count, "level2": l2_count, "level3": True, "errors": errors}
    top_level = [it for it in all_items if it.get("depth", 0) == 0]
    # If all items are at depth 0, use those; otherwise only top-level
    if not top_level:
        top_level = all_items

    l3_ok = False
    try:
        l3_ok = _run_meeting_briefing(
            meeting_id, top_level, all_items, client, mtg_model, briefing_prompt,
            max_tokens=mtg_max_tokens,
        )
    except Exception as exc:
        msg = f"Level 3 error: {exc}"
        logger.error(msg)
        errors.append(msg)

    db.set_meeting_status(meeting_id, "complete")
    _progress("Done.")

    return {"level1": l1_count, "level2": l2_count, "level3": l3_ok, "errors": errors}


# ---------------------------------------------------------------------------
# Ingest-time stub creators (called by pipeline/ingest.py to seed rows)
# ---------------------------------------------------------------------------

def summarize_agenda_item(item_id: int, created_by: str = "system") -> dict:
    """Create a placeholder summary_version row for an agenda item (status='stub')."""
    return db.create_summary_version(
        entity_type="agenda_item",
        entity_id=item_id,
        one_line=None,
        detailed=None,
        model_id=None,
        is_manual=False,
        status="stub",
        created_by=created_by,
    )


def summarize_meeting(meeting_id: int, created_by: str = "system") -> dict:
    """Create a placeholder summary_version row for a meeting (status='stub')."""
    return db.create_summary_version(
        entity_type="meeting",
        entity_id=meeting_id,
        one_line=None,
        detailed=None,
        model_id=None,
        is_manual=False,
        status="stub",
        created_by=created_by,
    )


def summarize_document(document_id: int, created_by: str = "system") -> dict:
    """Create a placeholder summary_version row for a document (status='stub')."""
    return db.create_summary_version(
        entity_type="document",
        entity_id=document_id,
        one_line=None,
        detailed=None,
        model_id=None,
        is_manual=False,
        status="stub",
        created_by=created_by,
    )


# ---------------------------------------------------------------------------
# Public read helpers (used by UI and tests)
# ---------------------------------------------------------------------------

def get_summary(entity_type: str, entity_id: int) -> dict | None:
    """Return the best current summary version for any entity."""
    return db.get_current_summary(entity_type, entity_id)


def list_summaries(entity_type: str, entity_id: int) -> list[dict]:
    """Return all summary versions for any entity, newest first."""
    return db.list_summary_versions(entity_type, entity_id)
